#!/usr/bin/env python3
"""Generate PPTX files from JSON slide deck definitions."""

import json
import sys
from pathlib import Path

def generate_pptx(slides_json, output_path):
    """Generate a PPTX file from slide definitions."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except ImportError:
        print("Error: python-pptx not installed. Run: pip install python-pptx")
        sys.exit(1)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for slide_def in slides_json:
        slide = prs.slides.add_slide(blank_layout)
        slide_type = slide_def.get("type", "content")
        title_text = slide_def.get("title", "")
        content = slide_def.get("content", [])

        if slide_type == "title":
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
            tf = txBox.text_frame
            tf.text = title_text
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(44)
                paragraph.font.bold = True
                paragraph.alignment = 1

        elif slide_type == "section":
            txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
            tf = txBox.text_frame
            tf.text = title_text
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(36)
                paragraph.font.bold = True
                paragraph.alignment = 1

        elif slide_type == "closing":
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
            tf = txBox.text_frame
            tf.text = title_text
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(40)
                paragraph.font.bold = True
                paragraph.alignment = 1
            if content:
                txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1.5))
                tf2 = txBox2.text_frame
                for item in content:
                    p = tf2.add_paragraph()
                    p.text = item
                    p.font.size = Pt(20)

        else:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
            tf = txBox.text_frame
            tf.text = title_text
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(32)
                paragraph.font.bold = True

            if content:
                txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(5))
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                for item in content:
                    p = tf2.add_paragraph()
                    p.text = f"• {item}"
                    p.font.size = Pt(20)
                    p.space_after = Pt(8)

    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: generate_pptx.py <slides.json> <output.pptx>")
        sys.exit(1)

    json_path = sys.argv[1]
    output_path = sys.argv[2]

    with open(json_path, "r", encoding="utf-8") as f:
        slides = json.load(f)

    result = generate_pptx(slides, output_path)
    print(f"Generated: {result}")
